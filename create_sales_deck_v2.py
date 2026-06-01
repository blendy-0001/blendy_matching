#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# カラーパレット定義
DARK = RGBColor(30, 39, 97)        # ネイビー
LIGHT = RGBColor(202, 220, 252)    # アイスブルー
ACCENT = RGBColor(255, 107, 107)   # コーラル
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(51, 51, 51)
GRAY = RGBColor(102, 102, 102)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_shape(slide, text, left, top, width, height, size=36, bold=True, color=DARK):
    """テキストボックスを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

    return textbox

def add_centered_text(slide, text, left, top, width, height, size=14, bold=False, color=TEXT):
    """センタリングされたテキストボックスを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    return textbox

def add_shape_with_text(slide, left, top, width, height, fill_color, text="", text_size=14, text_color=WHITE):
    """背景色付きの図形にテキストを追加"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.font.size = Pt(text_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    return shape

# ============================================
# プレゼンテーション初期化
# ============================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ============================================
# スライド 1: タイトルスライド
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK

add_centered_text(slide1, '協業マッチングシステム', 0.5, 1.8, 9, 1, size=48, bold=True, color=LIGHT)
add_centered_text(slide1, '最適なパートナーシップで、事業を拡大させる', 0.5, 2.7, 9, 0.6, size=20, color=LIGHT)
add_centered_text(slide1, 'Blendy Inc.', 0.5, 4.5, 9, 0.4, size=16, color=LIGHT)

# ============================================
# スライド 2: 問題提起
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide2, '現状の課題', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

issues = [
    {'title': '営業効率が低い', 'desc': '営業対象を見つけるのに時間がかかる'},
    {'title': 'リード獲得が困難', 'desc': '新規顧客へのアプローチ方法が限定的'},
    {'title': '無料だが手間がかかる', 'desc': '既存の無料マッチングサービスは自分たちで活用しないといけない'}
]

for idx, issue in enumerate(issues):
    y_pos = 1.2 + (idx * 1.15)

    # 背景ボックス
    shape = slide2.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(2)

    # タイトル
    add_title_shape(slide2, issue['title'], 0.8, y_pos + 0.1, 8.4, 0.35, size=16, bold=True, color=ACCENT)

    # 説明
    add_title_shape(slide2, issue['desc'], 0.8, y_pos + 0.45, 8.4, 0.4, size=12, bold=False, color=TEXT)

# ============================================
# スライド 3: ソリューション
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide3, 'Blendy の解決策', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

solutions = [
    {'title': 'マッチングまで実施', 'desc': '手をかけなくて良い。\nマッチングは弊社が実施'},
    {'title': 'AI 活用', 'desc': '複数の評価軸で\n最適なパートナーを\n自動選定'},
    {'title': '完全無料', 'desc': 'お客様負担なし。\n月2件送客させていただく\nビジネスモデル'}
]

for idx, sol in enumerate(solutions):
    x_pos = 0.5 + (idx * 3.2)

    # タイトル背景
    shape = slide3.shapes.add_shape(1, Inches(x_pos + 0.6), Inches(1.2), Inches(1), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = DARK
    shape.line.width = Pt(2)

    # アイコン
    add_centered_text(slide3, '●', x_pos + 0.6, 1.2, 1, 1, size=28, color=DARK)

    # タイトル
    add_centered_text(slide3, sol['title'], x_pos, 2.3, 3, 0.4, size=13, bold=True, color=DARK)

    # 説明
    add_centered_text(slide3, sol['desc'], x_pos, 2.8, 3, 1.2, size=10, bold=False, color=TEXT)

# ============================================
# スライド 4: 5つの相性軸
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide4, '協業を成功に導く 5つの相性軸', 0.5, 0.25, 9, 0.5, size=32, bold=True, color=DARK)
add_title_shape(slide4, '5つの軸が重なり合うことで、100点満点のスコアを算出', 0.5, 0.75, 9, 0.3, size=11, bold=False, color=GRAY)

# 5つの軸データ
axes = [
    {'num': '1', 'title': '顧客ターゲット\nの一致', 'desc': '業界・企業規模・\nターゲット層'},
    {'num': '2', 'title': '前後工程の\n補完性', 'desc': '職種パターンで\n相補性を判定'},
    {'num': '3', 'title': '市場での武器\n（情報）', 'desc': '80以上の業界\nキーワード判定'},
    {'num': '4', 'title': '事業のスケール\n設計', 'desc': '協業後の成長\n可能性を先読み'},
    {'num': '5', 'title': '市場方向性の\n合致', 'desc': '市場の方向性が\n一致しているか'}
]

# 中央から放射状に配置（5つの円）
center_x = 3.5
center_y = 2.5
radius = 1.3

import math

for idx, axis in enumerate(axes):
    # 角度計算（時計回りに配置）
    angle = (idx * 360 / 5 - 90) * math.pi / 180
    x_pos = center_x + radius * math.cos(angle)
    y_pos = center_y + radius * math.sin(angle)

    # 背景円
    shape = slide4.shapes.add_shape(1, Inches(x_pos - 0.35), Inches(y_pos - 0.35), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = DARK
    shape.line.width = Pt(2)

    # 番号
    add_centered_text(slide4, axis['num'], x_pos - 0.35, y_pos - 0.35, 0.7, 0.7, size=20, bold=True, color=WHITE)

    # タイトルボックス
    title_x = x_pos - 0.45 if x_pos < 3.5 else x_pos + 0.15 if x_pos > 3.5 else x_pos - 0.5
    title_y = y_pos - 0.6
    add_centered_text(slide4, axis['title'], title_x - 0.5, title_y, 1.0, 0.6, size=9, bold=True, color=DARK)

    # 説明
    desc_x = x_pos - 0.45 if x_pos < 3.5 else x_pos + 0.15 if x_pos > 3.5 else x_pos - 0.5
    desc_y = y_pos - 0.05
    add_centered_text(slide4, axis['desc'], desc_x - 0.5, desc_y, 1.0, 0.5, size=8, bold=False, color=GRAY)

# 中央に「100点満点」表示
center_shape = slide4.shapes.add_shape(1, Inches(center_x - 0.35), Inches(center_y - 0.35), Inches(0.7), Inches(0.7))
center_shape.fill.solid()
center_shape.fill.fore_color.rgb = DARK
center_shape.line.color.rgb = DARK
center_shape.line.width = Pt(2)

add_centered_text(slide4, '100\n点', center_x - 0.35, center_y - 0.35, 0.7, 0.7, size=14, bold=True, color=WHITE)

# ============================================
# スライド 5: ビジネスモデル
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide5, 'ビジネスモデル', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

# メイン訴求：完全無料
add_centered_text(slide5, '完全無料', 0.5, 0.9, 9, 0.5, size=44, bold=True, color=ACCENT)

# 説明セクション
model_items = [
    {'icon': '①', 'title': '無料登録', 'desc': '課題感をヒアリング'},
    {'icon': '②', 'title': '自動マッチング', 'desc': 'AI が相性の高いパートナーを選定'},
    {'icon': '③', 'title': '月数件送客', 'desc': '課題感に合うサービス提案を受け取る'},
]

y_pos = 1.7
for item in model_items:
    # アイコン
    icon_shape = slide5.shapes.add_shape(1, Inches(0.6), Inches(y_pos), Inches(0.5), Inches(0.5))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = ACCENT
    icon_shape.line.color.rgb = ACCENT
    add_centered_text(slide5, item['icon'], 0.6, y_pos, 0.5, 0.5, size=18, bold=True, color=WHITE)

    # タイトル
    add_title_shape(slide5, item['title'], 1.3, y_pos + 0.05, 2.5, 0.3, size=16, bold=True, color=DARK)

    # 説明
    add_title_shape(slide5, item['desc'], 1.3, y_pos + 0.35, 7.2, 0.3, size=12, bold=False, color=GRAY)

    y_pos += 0.8

# ボトムメッセージ
add_centered_text(slide5, '登録から成約まで、すべて完全無料でサポートします', 0.5, 4.8, 9, 0.5, size=13, bold=True, color=ACCENT)

# ============================================
# スライド 5.5: システムUI（ダッシュボード）
# ============================================
slide5_5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5_5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide5_5, 'かんたんな登録フローと自動マッチング', 0.5, 0.25, 9, 0.45, size=32, bold=True, color=DARK)

# ダッシュボードのイメージ説明
flow_items = [
    {'icon': '📋', 'title': 'ダッシュボード', 'desc': 'メンバー数、マッチング\n状況を一目で把握'},
    {'icon': '✍️', 'title': '複数活動登録', 'desc': '複数の事業・サービスを\n個別に登録可能'},
    {'icon': '🔄', 'title': '自動マッチング', 'desc': 'AI がリアルタイムで\n最適なペアを選定'}
]

for idx, item in enumerate(flow_items):
    x_pos = 0.6 + (idx * 3.1)

    # アイコン
    add_centered_text(slide5_5, item['icon'], x_pos, 1.2, 2.8, 0.6, size=32, bold=True)

    # タイトル
    add_centered_text(slide5_5, item['title'], x_pos, 1.9, 2.8, 0.35, size=13, bold=True, color=DARK)

    # 説明
    add_centered_text(slide5_5, item['desc'], x_pos, 2.3, 2.8, 0.9, size=10, bold=False, color=TEXT)

# 特徴
features = [
    '✓ 登録は5分で完了',
    '✓ 複数の事業・サービスごとに相性判定',
    '✓ リアルタイムでマッチング結果を確認'
]

feature_y = 3.4
for feature in features:
    add_title_shape(slide5_5, feature, 1.2, feature_y, 7.5, 0.3, size=11, bold=False, color=TEXT)
    feature_y += 0.35

# ============================================
# スライド 5: 導入実績
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide5, '導入実績', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

stats = [
    {'number': '13', 'label': '登録\nメンバー', 'color': DARK},
    {'number': '60', 'label': 'マッチング\n可能組数', 'color': ACCENT},
    {'number': '18', 'label': '累計マッチング\n完了', 'color': LIGHT}
]

for idx, stat in enumerate(stats):
    x_pos = 1.5 + (idx * 3)

    # 数字背景
    add_shape_with_text(slide5, x_pos, 1.3, 2.5, 2.5, stat['color'], stat['number'], text_size=52, text_color=WHITE if stat['color'] != LIGHT else DARK)

    # ラベル
    add_centered_text(slide5, stat['label'], x_pos, 2.9, 2.5, 0.6, size=12, bold=True, color=stat['color'] if stat['color'] != LIGHT else DARK)

# ============================================
# スライド 6: 利用方法（簡略化）
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide6, 'ご参加の流れ', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

simple_steps = [
    {'num': '1', 'title': '登録', 'desc': '企業情報と課題感を\nお聞きします'},
    {'num': '2', 'title': 'AI マッチング', 'desc': '最適なパートナーが\n自動選定されます'},
    {'num': '3', 'title': '成約サポート', 'desc': '契約までお手伝い\nさせていただきます'}
]

for idx, step in enumerate(simple_steps):
    x_pos = 0.8 + (idx * 3)

    # ステップ番号円
    add_shape_with_text(slide6, x_pos + 0.9, 1.3, 0.6, 0.6, ACCENT, step['num'], text_size=28)

    # タイトル
    add_centered_text(slide6, step['title'], x_pos, 2.1, 2.6, 0.35, size=14, bold=True, color=DARK)

    # 説明
    add_centered_text(slide6, step['desc'], x_pos, 2.5, 2.6, 1, size=11, bold=False, color=TEXT)

    # 矢印
    if idx < 2:
        add_centered_text(slide6, '→', x_pos + 2.7, 1.5, 0.3, 0.3, size=18, color=ACCENT)

# ============================================
# スライド 7: まとめ
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide7.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK

add_centered_text(slide7, 'Blendy で実現する', 0.5, 1.3, 9, 0.6, size=40, bold=True, color=LIGHT)
add_centered_text(slide7, 'リスクゼロの営業パートナーシップ', 0.5, 2.0, 9, 0.6, size=40, bold=True, color=WHITE)
add_centered_text(slide7, '完全無料 × 確実な相性マッチング', 0.5, 2.8, 9, 0.5, size=22, bold=True, color=ACCENT)

add_centered_text(slide7, 'お気軽にお問い合わせください', 0.5, 4.0, 9, 0.3, size=14, color=LIGHT)
add_centered_text(slide7, 'contact@blendyinc.com', 0.5, 4.4, 9, 0.35, size=13, bold=False, color=WHITE)

# ============================================
# 保存
# ============================================
output_path = r'C:\Users\yo_yo\Documents\blendy_matching\blendy_sales_deck.pptx'
prs.save(output_path)

print('Created sales deck: ' + output_path)
