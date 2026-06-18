#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# カラーパレット定義
DARK = RGBColor(30, 39, 97)        # ネイビー
LIGHT = RGBColor(202, 220, 252)    # アイスブルー
ACCENT = RGBColor(255, 107, 107)   # コーラル
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(51, 51, 51)
GRAY = RGBColor(102, 102, 102)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_shape(slide, text, left, top, width, height, size=36, bold=True, color=DARK):
    """テキストボックスを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

    return textbox

def add_centered_text(slide, text, left, top, width, height, size=14, bold=False, color=TEXT):
    """センタリングされたテキストボックスを追加"""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = text_frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    return textbox

def add_shape_with_text(slide, left, top, width, height, fill_color, text="", text_size=14, text_color=WHITE):
    """背景色付きの図形にテキストを追加"""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    if text:
        text_frame = shape.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.font.size = Pt(text_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    return shape

# ============================================
# プレゼンテーション初期化
# ============================================
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ============================================
# スライド 1: タイトルスライド
# ============================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白レイアウト
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK

add_centered_text(slide1, '協業マッチングシステム', 0.5, 1.8, 9, 1, size=48, bold=True, color=LIGHT)
add_centered_text(slide1, '最適なパートナーシップで、事業を拡大させる', 0.5, 2.7, 9, 0.6, size=20, color=LIGHT)
add_centered_text(slide1, 'Blendy Inc.', 0.5, 4.5, 9, 0.4, size=16, color=LIGHT)

# ============================================
# スライド 2: 問題提起
# ============================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide2, '現状の課題', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

issues = [
    {'title': '営業効率が低い', 'desc': '営業対象を見つけるのに時間がかかる'},
    {'title': 'リード獲得が困難', 'desc': '新規顧客へのアプローチ方法が限定的'},
    {'title': '相性の悪いマッチング', 'desc': '適切でないパートナーとの提携により失敗のリスク'}
]

for idx, issue in enumerate(issues):
    y_pos = 1.2 + (idx * 1.15)

    # 背景ボックス
    shape = slide2.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(2)

    # タイトル
    add_title_shape(slide2, issue['title'], 0.8, y_pos + 0.1, 8.4, 0.35, size=16, bold=True, color=ACCENT)

    # 説明
    add_title_shape(slide2, issue['desc'], 0.8, y_pos + 0.45, 8.4, 0.4, size=12, bold=False, color=TEXT)

# ============================================
# スライド 3: ソリューション
# ============================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide3, 'Blendy の解決策', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

solutions = [
    {'title': 'AI マッチング', 'desc': '複数の評価軸で最適なパートナーを自動選定'},
    {'title': '時間削減', 'desc': '営業活動の準備時間を 80% 削減'},
    {'title': '成功率向上', 'desc': '相性スコアに基づいた信頼できるマッチング'}
]

for idx, sol in enumerate(solutions):
    x_pos = 0.5 + (idx * 3.2)

    # タイトル背景
    shape = slide3.shapes.add_shape(1, Inches(x_pos + 0.6), Inches(1.2), Inches(1), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = DARK
    shape.line.width = Pt(2)

    # アイコン
    add_centered_text(slide3, '●', x_pos + 0.6, 1.2, 1, 1, size=28, color=DARK)

    # タイトル
    add_centered_text(slide3, sol['title'], x_pos, 2.35, 3, 0.4, size=14, bold=True, color=DARK)

    # 説明
    add_centered_text(slide3, sol['desc'], x_pos, 2.8, 3, 1.2, size=11, bold=False, color=TEXT)

# ============================================
# スライド 4: ビジネスモデル
# ============================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide4, 'ビジネスモデル', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)
add_title_shape(slide4, '成果報酬型（紹介される側のみ支払い）', 0.5, 1.0, 9, 0.4, size=14, bold=True, color=ACCENT)

pricing_items = [
    {'label': '性格マッチング', 'price': '¥5,000', 'desc': 'スキルセット相互補完なし'},
    {'label': '業種 + 性格マッチング', 'price': '¥10,000', 'desc': 'スキルセット相互補完あり'},
    {'label': '1商談（成約可能性高）', 'price': '¥25,000', 'desc': '3組セット'}
]

y_pos = 1.6
for item in pricing_items:
    # 背景ボックス
    shape = slide4.shapes.add_shape(1, Inches(0.5), Inches(y_pos), Inches(9), Inches(0.65))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 244, 255)
    shape.line.color.rgb = LIGHT
    shape.line.width = Pt(1)

    # ラベル
    add_title_shape(slide4, item['label'], 0.8, y_pos + 0.08, 4, 0.25, size=13, bold=True, color=DARK)

    # 価格
    tb = slide4.shapes.add_textbox(Inches(5.5), Inches(y_pos + 0.08), Inches(3.5), Inches(0.25))
    text_frame = tb.text_frame
    text_frame.text = item['price']
    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.RIGHT

    # 説明
    add_title_shape(slide4, item['desc'], 0.8, y_pos + 0.38, 8.2, 0.22, size=10, bold=False, color=GRAY)

    y_pos += 0.75

# ============================================
# スライド 5: 導入実績
# ============================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide5, '導入実績', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

stats = [
    {'number': '13', 'label': '登録\nメンバー', 'color': DARK},
    {'number': '60', 'label': 'マッチング\n可能組数', 'color': ACCENT},
    {'number': '18', 'label': '累計マッチング\n完了', 'color': LIGHT}
]

for idx, stat in enumerate(stats):
    x_pos = 1.5 + (idx * 3)

    # 数字背景
    add_shape_with_text(slide5, x_pos, 1.3, 2.5, 2.5, stat['color'], stat['number'], text_size=52, text_color=WHITE if stat['color'] != LIGHT else DARK)

    # ラベル
    add_centered_text(slide5, stat['label'], x_pos, 2.9, 2.5, 0.6, size=12, bold=True, color=stat['color'] if stat['color'] != LIGHT else DARK)

# ============================================
# スライド 6: 利用方法
# ============================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

add_title_shape(slide6, '簡単 3ステップ', 0.5, 0.3, 9, 0.5, size=36, bold=True, color=DARK)

steps = [
    {'num': '1', 'title': '登録', 'desc': 'フォームに\n企業情報を入力'},
    {'num': '2', 'title': 'マッチング実行', 'desc': 'AI が最適な\nパートナーを自動選定'},
    {'num': '3', 'title': '紹介・成約', 'desc': 'Blendy が\n契約までサポート'}
]

for idx, step in enumerate(steps):
    x_pos = 0.8 + (idx * 3)

    # ステップ番号円
    add_shape_with_text(slide6, x_pos + 0.9, 1.3, 0.6, 0.6, ACCENT, step['num'], text_size=28)

    # タイトル
    add_centered_text(slide6, step['title'], x_pos, 2.1, 2.6, 0.35, size=14, bold=True, color=DARK)

    # 説明
    add_centered_text(slide6, step['desc'], x_pos, 2.5, 2.6, 1, size=11, bold=False, color=TEXT)

    # 矢印
    if idx < 2:
        add_centered_text(slide6, '→', x_pos + 2.7, 1.5, 0.3, 0.3, size=18, color=ACCENT)

# ============================================
# スライド 7: まとめ
# ============================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide7.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK

add_centered_text(slide7, 'Blendy で実現する', 0.5, 1.5, 9, 0.6, size=40, bold=True, color=LIGHT)
add_centered_text(slide7, '最適なパートナーシップ', 0.5, 2.2, 9, 0.6, size=40, bold=True, color=WHITE)
add_centered_text(slide7, 'お気軽にお問い合わせください', 0.5, 3.5, 9, 0.4, size=16, color=LIGHT)
add_centered_text(slide7, 'contact@blendyinc.com', 0.5, 4.0, 9, 0.4, size=14, bold=False, color=WHITE)

# ============================================
# 保存
# ============================================
output_path = r'C:\Users\yo_yo\Documents\blendy_matching\営業資料_たたき.pptx'
prs.save(output_path)

print('Created sales deck: ' + output_path)
